#!/usr/bin/env python3
"""상영용 파워포인트 파일을 만듭니다.

slides.html 과 같은 내용, 같은 색을 씁니다.
호텔에 USB로 넘겨야 하거나 인터넷이 불안할 때 쓰세요.

    python3 make-pptx.py

두 벌이 나옵니다.
    루나-첫생일-동탄.pptx
    루나-첫생일-강릉.pptx
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ── 초대장과 같은 색 ──
PAPER = RGBColor(0xFC, 0xF8, 0xF1)
PAPER2 = RGBColor(0xF4, 0xED, 0xE1)
INK = RGBColor(0x33, 0x2E, 0x2A)
INK_SOFT = RGBColor(0x6B, 0x61, 0x55)
INK_FAINT = RGBColor(0x8A, 0x7F, 0x70)
LINE = RGBColor(0xE2, 0xD8, 0xC8)
SEAL = RGBColor(0xB9, 0x3A, 0x50)
SEAL_INK = RGBColor(0xFF, 0xF6, 0xF2)
BAND = [RGBColor(0xD9, 0x70, 0x7F), RGBColor(0xE8, 0xB3, 0x6B), RGBColor(0x8F, 0xB9, 0x8A),
        RGBColor(0x7F, 0xA3, 0xC4), RGBColor(0xA9, 0x8B, 0xB8)]

# 파워포인트에 없을 수 있는 글꼴은 피하고, 어디서나 있는 명조/고딕을 씁니다
SERIF = "맑은 고딕"          # 아래에서 제목만 바탕으로 바꿉니다
DISPLAY = "바탕"
BODY = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)   # 16:9

# 지나온 열두 달 영상. 파일이 없으면 안내 문구만 들어갑니다.
FILM = Path("video/diary_39s.mp4")
POSTER = Path("images/m-00.jpg")      # 재생 전에 보일 그림
THINK = Path("images/dolzabi-think.png")   # 돌잡이 결과 슬라이드에 쓰는 누끼 사진

BABY_FACTS = [
    ("태어난 날", "2025년 9월 22일"),
    ("좋아하는 것", "산책, 딸기, 아빠 안경"),
    ("요즘 하는 말", "엄마, 압빠, 맘마"),
]

GRABS = [("호랑이", "용기"), ("코끼리", "지혜"), ("강아지", "사랑"),
         ("나비", "희망"), ("말", "탐험"), ("거북이", "꾸준함")]

VENUES = {
    "동탄": dict(
        when="2026년 9월 12일 토요일",
        place="루나네 집",
        meal_title="이제 식사하러 갑니다",
        meal_lines=["오후 1시 · 긴자 신영통점", "경기 수원시 영통구 봉영로 1377"],
    ),
    "강릉": dict(
        when="2026년 9월 24일 목요일",
        place="씨마크 호텔 더 레스토랑",
        meal_title="이제 식사를 시작합니다",
        meal_lines=["편히 앉으셔서 맛있게 드세요"],
    ),
}


def blank(prs):
    """빈 슬라이드에 아이보리 배경을 깔아 돌려줍니다."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)          # 1 = 사각형
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, s, top, height, size, color, font=BODY, bold=False,
         spacing=0.0, align=PP_ALIGN.CENTER, line=1.4, left=None, width=None):
    box = slide.shapes.add_textbox(left if left is not None else Inches(0.8),
                                   top,
                                   width if width is not None else W - Inches(1.6),
                                   height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
        r.font.bold = bold
        if spacing:
            # 자간은 python-pptx가 직접 지원하지 않아 XML로 넣습니다
            r.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def rule(slide, top, height=Inches(0.42)):
    r = slide.shapes.add_shape(1, W // 2 - Emu(6350), top, Emu(12700), height)
    r.fill.solid()
    r.fill.fore_color.rgb = LINE
    r.line.fill.background()
    r.shadow.inherit = False


def band(slide, top):
    seg_w, gap = Inches(0.5), Inches(0.05)
    total = seg_w * 5 + gap * 4
    x = W // 2 - total // 2
    for c in BAND:
        b = slide.shapes.add_shape(1, x, top, seg_w, Emu(38100))
        b.fill.solid()
        b.fill.fore_color.rgb = c
        b.line.fill.background()
        b.shadow.inherit = False
        x += seg_w + gap


def seal(slide, top):
    size = Inches(1.05)
    sh = slide.shapes.add_shape(1, W // 2 - size // 2, top, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = SEAL
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.rotation = -7
    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "돌"
    r.font.size = Pt(30)
    r.font.color.rgb = SEAL_INK
    r.font.name = DISPLAY


def build(venue_name, v, out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 표지
    s = blank(prs)
    seal(s, Inches(0.85))
    text(s, "초 대 합 니 다", Inches(2.15), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "김 루 나", Inches(2.7), Inches(1.5), 66, INK, DISPLAY, spacing=3)
    rule(s, Inches(4.35))
    text(s, "첫 번째 생일", Inches(4.9), Inches(0.6), 24, INK_SOFT, DISPLAY)
    text(s, f"{v['when']} · {v['place']}", Inches(5.7), Inches(0.5), 15, INK_FAINT, BODY)

    # 2 환영
    s = blank(prs)
    text(s, "환 영 합 니 다", Inches(1.6), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "와 주셔서\n고맙습니다", Inches(2.3), Inches(2.4), 58, INK, DISPLAY, line=1.3)
    band(s, Inches(5.4))

    # 3 루나는요
    s = blank(prs)
    text(s, "우 리  아 이", Inches(0.9), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "루나는요", Inches(1.5), Inches(1.1), 50, INK, DISPLAY)
    top = Inches(3.1)
    for k, val in BABY_FACTS:
        text(s, k, top, Inches(0.6), 15, INK_FAINT, BODY, align=PP_ALIGN.RIGHT,
             left=Inches(2.6), width=Inches(3.0))
        text(s, val, top, Inches(0.6), 27, INK, DISPLAY, align=PP_ALIGN.LEFT,
             left=Inches(6.1), width=Inches(5.5))
        top += Inches(0.95)

    # 4 영상
    s = blank(prs)
    text(s, "사 진", Inches(0.42), Inches(0.45), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "지나온 열두 달", Inches(0.95), Inches(0.85), 36, INK, DISPLAY)
    # 16:9 영상을 화면 가운데에 놓습니다
    vid_w = Inches(8.4)
    vid_h = Inches(8.4 * 9 / 16)
    if FILM.exists():
        s.shapes.add_movie(str(FILM), W // 2 - vid_w // 2, Inches(2.05), vid_w, vid_h,
                           poster_frame_image=str(POSTER) if POSTER.exists() else None,
                           mime_type="video/mp4")
    else:
        text(s, "— 여기서 영상을 틀어 주세요 —", Inches(3.6), Inches(0.5), 16, INK_FAINT, BODY)
    text(s, "재생 버튼을 누르면 시작합니다", Inches(6.85), Inches(0.45), 13, INK_FAINT, BODY)

    # 5 돌잡이 열기
    s = blank(prs)
    text(s, "돌 잡 이", Inches(0.9), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "어떤 사람으로\n자랄까요?", Inches(1.5), Inches(2.3), 52, INK, DISPLAY, line=1.3)
    rule(s, Inches(4.25))
    text(s, "무엇이 될지보다\n어떤 마음으로 자랄지가 궁금했습니다",
         Inches(4.9), Inches(1.5), 24, INK_SOFT, DISPLAY, line=1.8)

    # 6 여섯 동물
    s = blank(prs)
    text(s, "돌 잡 이", Inches(0.55), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "여섯 가지 성품", Inches(1.1), Inches(0.9), 40, INK, DISPLAY)
    cw, ch, gap = Inches(3.5), Inches(1.85), Inches(0.28)
    grid_w = cw * 3 + gap * 2
    x0 = W // 2 - grid_w // 2
    y0 = Inches(2.5)
    for i, (animal, meaning) in enumerate(GRABS):
        cx = x0 + (i % 3) * (cw + gap)
        cy = y0 + (i // 3) * (ch + gap)
        card = s.shapes.add_shape(1, cx, cy, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = PAPER2
        card.line.color.rgb = LINE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = animal
        r.font.size = Pt(32)
        r.font.color.rgb = INK
        r.font.name = DISPLAY
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = meaning
        r2.font.size = Pt(16)
        r2.font.color.rgb = SEAL
        r2.font.name = BODY

    # 7 결과
    s = blank(prs)
    text(s, "돌 잡 이", Inches(0.5), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "루나가 잡은 것은", Inches(1.05), Inches(1.15), 48, INK, DISPLAY)
    if THINK.exists():
        # 배경을 지운 사진이라 아이보리 위에 그대로 올려 놓습니다
        th_h = Inches(4.15)
        th_w = Emu(int(th_h * 820 / 893))
        s.shapes.add_picture(str(THINK), W // 2 - th_w // 2, Inches(2.35), th_w, th_h)
        band(s, Inches(6.85))
    else:
        band(s, Inches(5.2))

    # 8 맞히신 분
    s = blank(prs)
    text(s, "돌 잡 이", Inches(0.95), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "맞히신 분\n계신가요?", Inches(1.5), Inches(2.2), 50, INK, DISPLAY, line=1.3)
    rule(s, Inches(4.0))
    text(s, "초대장을 다시 열어 보시면\n고르셨던 동물이 그대로 남아 있습니다",
         Inches(4.6), Inches(1.5), 24, INK_SOFT, DISPLAY, line=1.8)
    text(s, "화면을 보여 주세요", Inches(6.3), Inches(0.5), 15, INK_FAINT, BODY)

    # 9 노래
    s = blank(prs)
    text(s, "축 하  노 래", Inches(0.9), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "생일 축하합니다\n생일 축하합니다\n사랑하는 우리 루나\n생일 축하합니다",
         Inches(1.8), Inches(4.6), 40, INK, DISPLAY, line=1.75)

    # 10 숫자 뽑기
    s = blank(prs)
    text(s, "행 운 의  숫 자", Inches(0.95), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "번호를 뽑아 주세요", Inches(1.6), Inches(1.2), 48, INK, DISPLAY)
    rule(s, Inches(3.15))
    text(s, "가장 작은 수와 가장 큰 수를 뽑으신 분께\n작은 선물을 드립니다",
         Inches(3.8), Inches(1.6), 25, INK_SOFT, DISPLAY, line=1.8)
    text(s, "초대장 맨 아래 ‘첫돌’ 글자를 세 번 누르세요",
         Inches(5.7), Inches(0.6), 20, SEAL, DISPLAY)

    # 11 사진
    s = blank(prs)
    text(s, "사 진", Inches(1.0), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "다 같이\n사진 찍겠습니다", Inches(2.2), Inches(2.4), 52, INK, DISPLAY, line=1.3)
    band(s, Inches(5.4))

    # 12 감사
    s = blank(prs)
    text(s, "고 맙 습 니 다", Inches(1.0), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, "덕분에\n한 해를 잘 지냈습니다", Inches(1.7), Inches(2.4), 48, INK, DISPLAY, line=1.35)
    rule(s, Inches(4.6))
    text(s, "오래 기억할 하루가 되겠습니다", Inches(5.2), Inches(0.7), 25, INK_SOFT, DISPLAY)

    # 13 식사
    s = blank(prs)
    text(s, "식 사", Inches(1.1), Inches(0.5), 15, INK_FAINT, DISPLAY, spacing=4)
    text(s, v["meal_title"], Inches(1.8), Inches(1.2), 48, INK, DISPLAY)
    rule(s, Inches(3.4))
    text(s, "\n".join(v["meal_lines"]), Inches(4.0), Inches(1.5), 26, INK_SOFT, DISPLAY, line=1.8)
    band(s, Inches(5.9))

    prs.save(out)
    return len(prs.slides.__iter__.__self__._sldIdLst)


def main() -> int:
    for name, v in VENUES.items():
        out = f"루나-첫생일-{name}.pptx"
        n = build(name, v, out)
        print(f"{out} — {n}장")
    return 0


if __name__ == "__main__":
    sys.exit(main())
